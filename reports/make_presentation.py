"""
Generate the Sonder product presentation as a PPTX file.

Visual identity matches the website tokens:
  - Dark velvet background (#080807)
  - Gold accents (#D4B686)
  - Cream / bone text (#F4EDE0)
  - Serif italic headlines, sans-serif body
  - Gold pulse dots + tracked-out uppercase eyebrows
  - Stylised gold S mark on title and closing slides

Output: reports/sonder_presentation.pptx
Run:    python reports/make_presentation.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt


# Brand tokens (mirror jahnvi/frontend/src/lib/tokens.js).
BG       = RGBColor(0x08, 0x08, 0x07)
GOLD     = RGBColor(0xD4, 0xB6, 0x86)
BONE     = RGBColor(0xF4, 0xED, 0xE0)
MUTE     = RGBColor(0xA9, 0xA0, 0x8D)
DIM      = RGBColor(0x6F, 0x68, 0x5A)
HAIRLINE = RGBColor(0x38, 0x33, 0x29)

SERIF_ITALIC = "Georgia"      # PPT fallback for Cormorant Garamond italic
SANS         = "Helvetica"    # PPT fallback for Inter Tight

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN  = Inches(0.6)


def apply_background(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def add_text(slide, x, y, w, h, text, *,
             font=SERIF_ITALIC, size=24, italic=True,
             color=BONE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             letter_spacing=None, line_spacing=1.15, bold=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if letter_spacing is not None:
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", str(letter_spacing))
    return box


def add_eyebrow(slide, x, y, w, text, color=GOLD):
    return add_text(slide, x, y, w, Inches(0.3), text.upper(),
                    font=SANS, size=9, italic=False, color=color,
                    letter_spacing=320)


def add_pulse_dot(slide, x, y, color=GOLD):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.08), Inches(0.08))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    return dot


def add_gold_underline(slide, x, y, w=Inches(1.5)):
    add_rect(slide, x, y, w, Emu(12700), GOLD)


def add_sonder_mark(slide, cx_in, cy_in, size_in=1.2):
    """Stylised gold S — two angled blades forming an S silhouette,
    next to the Sonder wordmark in serif italic."""
    half = size_in / 2
    cx = Inches(cx_in)
    cy = Inches(cy_in)

    # Top blade (rotated 180 so it angles top-right → bottom-left).
    top = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE,
        cx - Inches(half * 0.45), cy - Inches(half * 0.95),
        Inches(size_in * 0.85), Inches(size_in * 0.55),
    )
    top.rotation = 180

    # Bottom blade (default orientation, mirrors top across the centre).
    bot = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE,
        cx - Inches(half * 0.40), cy + Inches(half * 0.05),
        Inches(size_in * 0.85), Inches(size_in * 0.55),
    )

    for blade in (top, bot):
        blade.fill.solid()
        blade.fill.fore_color.rgb = GOLD
        blade.line.fill.background()

    # Wordmark
    add_text(slide,
             cx + Inches(size_in * 0.6), cy - Inches(size_in * 0.32),
             Inches(3.5), Inches(0.6),
             "Sonder", font=SERIF_ITALIC, size=int(34 * size_in),
             italic=True, color=BONE, align=PP_ALIGN.LEFT)


def slide_blank(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, BG)
    return slide


def page_footer(slide, page_num, total):
    add_text(slide, MARGIN, SLIDE_H - Inches(0.5),
             Inches(3), Inches(0.3),
             "SONDER",
             font=SANS, size=8, italic=False, color=DIM, letter_spacing=320)
    add_text(slide, SLIDE_W - MARGIN - Inches(1.5), SLIDE_H - Inches(0.5),
             Inches(1.5), Inches(0.3),
             f"{page_num:02d} / {total:02d}",
             font=SANS, size=8, italic=False, color=DIM,
             letter_spacing=320, align=PP_ALIGN.RIGHT)


def title_slide(prs, total):
    s = slide_blank(prs)
    add_sonder_mark(s, cx_in=4.0, cy_in=3.0, size_in=1.4)
    add_text(s, MARGIN, Inches(4.0), SLIDE_W - 2 * MARGIN, Inches(1.5),
             "Trip planning, finally built for who you are.",
             font=SERIF_ITALIC, size=44, color=BONE)
    add_text(s, MARGIN, Inches(5.2), SLIDE_W - 2 * MARGIN, Inches(0.6),
             "A GenAI travel-planning and co-traveller-matching system.",
             font=SANS, size=14, italic=False, color=MUTE)
    add_gold_underline(s, MARGIN, Inches(5.95))
    add_text(s, MARGIN, Inches(6.15), SLIDE_W - 2 * MARGIN, Inches(0.4),
             "PRODUCT OVERVIEW",
             font=SANS, size=10, italic=False, color=GOLD, letter_spacing=360)
    page_footer(s, 1, total)


def bullets_slide(prs, page_num, total, eyebrow, headline, bullets):
    s = slide_blank(prs)
    add_pulse_dot(s, MARGIN, Inches(0.95))
    add_eyebrow(s, MARGIN + Inches(0.18), Inches(0.92), Inches(6), eyebrow)
    add_text(s, MARGIN, Inches(1.4), SLIDE_W - 2 * MARGIN, Inches(1.4),
             headline,
             font=SERIF_ITALIC, size=36, color=BONE, line_spacing=1.1)
    add_gold_underline(s, MARGIN, Inches(2.7), w=Inches(2))
    cur_y = Inches(3.05)
    for label, body in bullets:
        add_rect(s, MARGIN, cur_y + Inches(0.15), Inches(0.06), Inches(0.06), GOLD)
        add_text(s, MARGIN + Inches(0.22), cur_y, Inches(2.8), Inches(0.4),
                 label,
                 font=SERIF_ITALIC, size=17, color=GOLD)
        add_text(s, MARGIN + Inches(3.1), cur_y, SLIDE_W - MARGIN - Inches(3.5), Inches(0.8),
                 body,
                 font=SANS, size=12, italic=False, color=BONE, line_spacing=1.4)
        cur_y += Inches(0.78)
    page_footer(s, page_num, total)


def closing_slide(prs, page_num, total):
    s = slide_blank(prs)
    add_sonder_mark(s, cx_in=4.0, cy_in=3.0, size_in=1.2)
    add_text(s, MARGIN, Inches(4.2), SLIDE_W - 2 * MARGIN, Inches(0.9),
             "Travel, together.",
             font=SERIF_ITALIC, size=40, color=BONE)
    add_gold_underline(s, MARGIN, Inches(5.1), w=Inches(2.5))
    add_text(s, MARGIN, Inches(5.3), SLIDE_W - 2 * MARGIN, Inches(0.4),
             "BUILT BY A TEAM TIRED OF TRIP-PLANNING ALGORITHMS",
             font=SANS, size=10, italic=False, color=GOLD, letter_spacing=360)
    add_text(s, MARGIN, Inches(5.8), SLIDE_W - 2 * MARGIN, Inches(0.4),
             "github.com/shreyast36/sonder",
             font=SANS, size=12, italic=False, color=MUTE)
    page_footer(s, page_num, total)


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    TOTAL = 19
    n = 1

    title_slide(prs, TOTAL); n += 1

    bullets_slide(prs, n, TOTAL,
        "The problem",
        "Every existing platform optimises for the wrong thing.",
        [
            ("TripAdvisor",   "Surfaces what coach tours stop at. Top three things in any city become the same three things on every list."),
            ("Google Travel", "Knows your budget but not your mood."),
            ("Instagram",     "Knows your mood but not whether the place is closed on Tuesdays."),
            ("Airbnb",        "Solves where you sleep and stops there."),
            ("All of them",   "Nobody solves the harder problem of who you travel with."),
        ]); n += 1

    bullets_slide(prs, n, TOTAL,
        "What Sonder is",
        "Three things at once that nobody else does together.",
        [
            ("Itineraries",   "Day-by-day plans that name real places — the café with the rosé awnings on Largo do Carmo, not 'explore the old town'."),
            ("Matching",      "Solo and couple travellers paired against a measurable compatibility surface — not a free-text bio."),
            ("Social layer",  "192 LLM-designed synthetic travellers post, open trips, and message real users from day one."),
        ]); n += 1

    bullets_slide(prs, n, TOTAL,
        "Architecture",
        "Four cleanly separated service planes.",
        [
            ("Frontend",       "React + Vite served from Cloudflare's global edge. Catch-all Pages Function proxies every /api/* to the backend."),
            ("Backend",        "FastAPI on Render. Single process handles REST, SSE, WebSockets, and the synthetic-agents loop in one lifespan."),
            ("Language models","Anthropic Claude primary, OpenAI fallback. Small tier for fast persona voice, large tier for deep generation, validator tier on either side."),
            ("Data stores",    "Pinecone for vectors, Firestore for operational state, Firebase Storage for binary assets, all split by access pattern."),
        ]); n += 1

    bullets_slide(prs, n, TOTAL,
        "Foundations",
        "Three frameworks the system is built on.",
        [
            ("PPM theory",  "Dann (1977), Crompton (1979). 12 dimensions of travel motivation — six push, six pull. Auditable substring keyword matching."),
            ("GoEmotions",  "Demszky et al. (2020). 27 emotion labels used as anchor vectors in our embedding space. Defensible cosine classification."),
            ("Signature",   "Eight-key emotional taxonomy synthesised by the team. Drives private framing for every persona-voiced surface."),
        ]); n += 1

    bullets_slide(prs, n, TOTAL,
        "Vector database",
        "One Pinecone index, three namespaces, unified embedding space.",
        [
            ("Destinations",   "Hand-curated ~50 cities. Embedded text describes what each place feels like, not just facts."),
            ("Activities",     "~500 real venues, per-destination. Cosine retrieval anchored on the user's persona and emotional signature."),
            ("Cotravellers",   "192 LLM-designed synthetic personas + 18 couples. Same vector space so cross-namespace queries are coherent."),
            ("Embedding",      "OpenAI text-embedding-3-small, 1536 dimensions. One model means all corpora share the same geometry."),
        ]); n += 1

    bullets_slide(prs, n, TOTAL,
        "Model stack",
        "Eight distinct models, each chosen for its task profile.",
        [
            ("Claude Haiku 4.5",       "Small-tier conversational surfaces — chat, openers, classifiers, why-this explanations."),
            ("Claude Sonnet 4.6",      "Large-tier generative surfaces — itinerary generation, complex refinement."),
            ("OpenAI GPT-4o family",   "Per-tier fallback when Anthropic is unavailable. Per-provider model ids prevent cross-provider misroutes."),
            ("text-embedding-3-small", "All retrieval embeddings, 1536-dim."),
            ("gpt-image-1, ElevenLabs","Synthetic-persona portraits at seed time; persona voice text-to-speech for chat playback."),
        ]); n += 1

    bullets_slide(prs, n, TOTAL,
        "Itinerary generation",
        "Retrieval-grounded, validator-gated, streamed day-by-day.",
        [
            ("Information starvation","Three persona pipelines run in parallel — embedder, GoEmotions, dimension labeller — each blind to the others."),
            ("Hard pre-filters",      "Budget feasibility, avoid-list veto, must-haves all enforced before re-ranking."),
            ("Streaming UX",          "Forward-only JSON parser yields each day the moment its brace closes — user sees Day 1 in 15s while Day 7 is still generating."),
            ("Validator loop",        "Seven critic categories. Failed outputs go through up-to-three refinement attempts with fresh embeddings."),
        ]); n += 1

    bullets_slide(prs, n, TOTAL,
        "Co-traveller matching",
        "A calibrated probability — not a uniform cosine.",
        [
            ("Feature pipeline",     "Ten reusable scoring functions across cosine, persona overlap, emotional signature, pace, budget, interests."),
            ("Hard safety filters",  "Travel-style and same-gender (solo) enforced before re-ranking. Couples never see solo personas."),
            ("Score-as-probability", "Score 0.6 honestly means roughly 60% chance the matched persona accepts. Denial becomes meaningful."),
            ("Equal-weight priors",  "1/N across features. Honest about uncertainty; every shown candidate is logged for phase-two gradient learning."),
        ]); n += 1

    bullets_slide(prs, n, TOTAL,
        "Validator engine",
        "Five critic prompts that keep personas in character.",
        [
            ("Assistant-voice leakage","'How can I help you?' caught before it reaches the user."),
            ("AI / tooling leakage",   "'As an AI', 'I'm a language model' blocked."),
            ("Memory contradictions",  "Regex on chat history catches personas claiming trips they just denied."),
            ("Token-level failures",   "Empty generations, repetition stutters, malformed JSON."),
            ("Taxonomy leakage",       "Persona uses 'push' or 'pull' as words — caught and rewritten."),
        ]); n += 1

    bullets_slide(prs, n, TOTAL,
        "Learning loop",
        "User feedback applies weight penalties; the system re-runs.",
        [
            ("Explicit feedback",   "'Cheaper' maps to budget-fit features. 'Less packed' maps to pace. Decay across turns prevents oscillation."),
            ("Self-correction",     "The generation pipeline re-runs end-to-end with updated weights. The system literally re-prices its own previous recommendation."),
            ("Implicit chat signal","Sarcasm-aware, negation-aware scanner re-ranks the candidate after every chat message."),
            ("Auditable",           "Every revision turn records feedback, weight delta, dropped + added titles, validator verdict."),
        ]); n += 1

    bullets_slide(prs, n, TOTAL,
        "Real-time layer",
        "WebSockets for conversation; Firestore listeners for state.",
        [
            ("First-message auth",   "Tokens go in the first JSON message after handshake — never in query strings. Cleaner audit trail."),
            ("TTL-derived presence", "30-second pings; presence is (now − last_seen) < 90s, not a boolean flag that goes stale on disconnect."),
            ("Firestore listeners",  "Shared-itinerary edits propagate to both participants in milliseconds via the native subscription API."),
            ("Modular boundaries",   "Four team-owned folders with strict import boundaries. Schema changes don't break LLM calls."),
        ]); n += 1

    bullets_slide(prs, n, TOTAL,
        "Shared itinerary",
        "Two travellers plan one trip together, mediated by the persona.",
        [
            ("Propose-counter loop", "No hard reject — the persona must accept or counter. Negotiation always moves forward."),
            ("Reason-in-message",    "Every counterproposal carries a brief conversational reason. Counters without reasons are caught by the validator."),
            ("Optimistic locking",   "Every write checks client_version. 409 on mismatch; client re-fetches and retries."),
            ("Token-Jaccard dedupe", "Counters can't recycle previously-rejected activities. If the LLM tries, we flip the verdict to accept."),
        ]); n += 1

    bullets_slide(prs, n, TOTAL,
        "Synthetic population",
        "192 LLM-designed personas, blind-writer pipeline.",
        [
            ("Diversity matrix", "16 cities × 3 age buckets × 2 genders × 2 per cell. 50/50 gender split hard-locked."),
            ("Two-stage design", "Writer LLM gets only city, age, gender. Inferrer assigns PPM dimensions and emotional signature blind to the writer's output."),
            ("Stylised portraits","gpt-image-1 painterly outputs — explicitly not photorealistic, biased that way for the 'Sonder Curated' disclosure."),
            ("Backfill template","Schema-only updates patch metadata via index.update without re-paying LLM cost. Pattern reused for future fields."),
        ]); n += 1

    bullets_slide(prs, n, TOTAL,
        "Approve flow",
        "The trip you locked in deserves a moment.",
        [
            ("Ken-Burns montage", "Five destination photos cycle every 1.2 seconds, each crossfading with a slow zoom."),
            ("Typographic build", "LOCKED IN tracks out → COUNTRY drops in → CITY slams from 1.55x scale + 16px blur to 1x / blur 0."),
            ("Gold-dust particles","Drift top-to-bottom under a cinematic vignette. Trailer cadence, not browser tab."),
            ("Co-traveller prompt","After the reveal, asks if you want to plan with someone. Yes → matching; no → dashboard."),
        ]); n += 1

    bullets_slide(prs, n, TOTAL,
        "Group-style filtering",
        "Solo, couple, family, friends are four different products.",
        [
            ("Solo",   "Counter-seating venues, walking distance, 1-2 meet-others activities. Same-gender hard filter on matching."),
            ("Couple", "Every overnight private. Tables for two. At least one slow shared activity per day."),
            ("Family", "Assumes kids: kids menus, dinner by 7pm, mid-day reset, no clubs. Matching disabled — they have their party."),
            ("Friends","One-table reservations for N, split-and-rejoin activities, apartment over hotel rooms. Matching disabled."),
        ]); n += 1

    bullets_slide(prs, n, TOTAL,
        "Observability",
        "Sentry and PostHog — deliberately separate tools.",
        [
            ("Sentry",            "Errors only. Filter drops 'Failed to fetch' TypeErrors and expected 4xx so the feed stays signal-dense."),
            ("PostHog",           "Product analytics. Five top-level families: satisfaction, retrieval quality, response quality, hallucination, completion funnel."),
            ("Validator telemetry","First-try approval, repair count, semantic-genericity score. Continuous signal catches drift weeks early."),
            ("Feature stats",     "Per-feature distribution observer catches silent scale domination before it shows up as flat engagement."),
        ]); n += 1

    bullets_slide(prs, n, TOTAL,
        "Deployment",
        "Operational independence at scale.",
        [
            ("Cloudflare Pages",  "Global edge. SPA on a single domain with a catch-all Pages Function proxying /api/*."),
            ("Render",            "Long-running FastAPI container. Auto-deploy from main, TLS termination, health checks."),
            ("Multi-provider",    "Per-tier LLM provider preference with per-provider model ids — no cross-provider model id misroutes."),
            ("Failure isolation", "Cloudflare incident? Backend stays reachable. Anthropic rate-limited? OpenAI picks up. Every layer replaceable independently."),
        ]); n += 1

    closing_slide(prs, n, TOTAL)

    out = Path(__file__).resolve().parent / "sonder_presentation.pptx"
    prs.save(str(out))
    print(f"PPTX generated: {out}")
    return out


if __name__ == "__main__":
    build()
